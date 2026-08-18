import os
import json
from pypdf import PdfReader
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
import pytesseract
from PIL import Image
import speech_recognition as sr
from pydub import AudioSegment
from moviepy import VideoFileClip

# Check for Tesseract executable in common paths (Windows)
# Users might need to add it to PATH or set it here.
# A robust solution tries to find it.
possible_tesseract_paths = [
    r"C:\\Program Files\\Tesseract-OCR\\tesseract.exe",
    r"C:\\Program Files (x86)\\Tesseract-OCR\\tesseract.exe"
]
for path in possible_tesseract_paths:
    if os.path.exists(path):
        pytesseract.pytesseract.tesseract_cmd = path
        break

class ParserFactory:
    @staticmethod
    def get_content(file_path):
        """
        Detects file type and returns the extracted text content.
        """
        ext = os.path.splitext(file_path)[1].lower()
        
        try:
            if ext == '.json':
                return parse_json(file_path)
            elif ext == '.pdf':
                return parse_pdf(file_path)
            elif ext == '.txt':
                return parse_text(file_path)
            elif ext == '.md':
                return parse_text(file_path)
            elif ext in ['.docx', '.doc']:
                return parse_docx(file_path)
            elif ext == '.xlsx':
                return parse_xlsx(file_path)
            elif ext == '.pptx':
                return parse_pptx(file_path)
            elif ext in ['.jpg', '.jpeg', '.png', '.bmp', '.tiff']:
                return parse_image(file_path)
            elif ext in ['.mp3', '.wav', '.m4a', '.flac']:
                return parse_audio(file_path)
            elif ext in ['.mp4', '.avi', '.mov', '.mkv']:
                return parse_video(file_path)
            else:
                return f"[Unsupported file type: {ext}]"
        except Exception as e:
            return f"[Error parsing {os.path.basename(file_path)}: {str(e)}]"

def parse_json(path):
    with open(path, 'r', encoding='utf-8', errors='ignore') as f:
        data = json.load(f)
        # Attempt to find content fields
        if isinstance(data, dict):
            return data.get('content') or data.get('body') or data.get('text') or str(data)
        elif isinstance(data, list):
            return " ".join([str(item) for item in data])
        return str(data)

def parse_text(path):
    with open(path, 'r', encoding='utf-8', errors='ignore') as f:
        return f.read()

def parse_pdf(path):
    text = ""
    try:
        reader = PdfReader(path)
        for page in reader.pages:
            extracted = page.extract_text()
            if extracted:
                text += extracted + "\n"
    except Exception as e:
        return f"[PDF Error: {e}]"
    return text if text else "[PDF contained no readable text]"

def parse_docx(path):
    doc = Document(path)
    return "\n".join([para.text for para in doc.paragraphs])

def parse_xlsx(path):
    wb = load_workbook(path, data_only=True)
    text = []
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        text.append(f"--- Sheet: {sheet} ---")
        for row in ws.iter_rows(values_only=True):
            row_text = " ".join([str(cell) for cell in row if cell is not None])
            if row_text:
                text.append(row_text)
    return "\n".join(text)

def parse_pptx(path):
    prs = Presentation(path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

def parse_image(path):
    try:
        image = Image.open(path)
        text = pytesseract.image_to_string(image)
        return text if text.strip() else "[OCR found no text]"
    except pytesseract.TesseractNotFoundError:
        return "[Tesseract OCR not installed or not in PATH. Please install Tesseract-OCR]"
    except Exception as e:
        return f"[OCR Error: {e}]"

def parse_audio(path):
    # Convert to WAV for SpeechRecognition if needed
    try:
        r = sr.Recognizer()
        
        # Helper to convert
        wav_path = path
        converted = False
        if not path.endswith('.wav'):
            audio = AudioSegment.from_file(path)
            wav_path = "temp_audio.wav"
            audio.export(wav_path, format="wav")
            converted = True
            
        with sr.AudioFile(wav_path) as source:
            audio_data = r.record(source)
            text = r.recognize_google(audio_data) # Uses Google Web Speech API (Free tier limits) 
            
        if converted and os.path.exists(wav_path):
            os.remove(wav_path)
            
        return f"[Audio Transcript]: {text}"
    except sr.RequestError as e:
        return f"[Speech Recognition API unavailable: {e}]"
    except sr.UnknownValueError:
        return "[Audio unintelligible]"
    except Exception as e:
        return f"[Audio Processing Error: {e}. Note: FFmpeg is required for non-WAV files.]"

def parse_video(path):
    try:
        # Extract audio from video
        video = VideoFileClip(path)
        audio_path = "temp_video_audio.wav"
        video.audio.write_audiofile(audio_path, logger=None)
        
        # Use audio parser logic
        r = sr.Recognizer()
        with sr.AudioFile(audio_path) as source:
            audio_data = r.record(source)
            text = r.recognize_google(audio_data)
            
        # Clean up
        if os.path.exists(audio_path):
            os.remove(audio_path)
        video.close()
        
        return f"[Video Transcript]: {text}"
    except Exception as e:
        return f"[Video Error: {e}. Note: FFmpeg required.]"
